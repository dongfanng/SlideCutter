"""读取 layout.json，用 python-pptx 重建为 PPT。

layout.json 由前端切片工具导出，结构示例：

    {
      "source": "material.png",
      "width": 1200,
      "height": 800,
      "elements": [
        { "file": "element-01.png", "index": 0,
          "x": 10, "y": 20, "width": 100, "height": 120 }
      ]
    }

坐标单位为原图像素。默认按 96 DPI 换算为英寸（1 px = 1/96 in），
并在画布上按原位置摆放每个透明 PNG，从而复现原始排版。
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from collections.abc import Mapping
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# 像素 → 英寸 的换算（屏幕 96 DPI，1 英寸 = 96 像素）
PX_PER_INCH = 96
# 缺少画布宽高时的回退尺寸（16:9，单位英寸）
DEFAULT_SLIDE_W = 13.333
DEFAULT_SLIDE_H = 7.5


def load_layout(layout_path: Path) -> dict:
    """读取并校验 layout.json。"""
    if not layout_path.is_file():
        raise FileNotFoundError(f"layout.json 不存在：{layout_path}")
    with layout_path.open(encoding="utf-8") as f:
        layout = json.load(f)
    if not isinstance(layout, dict) or "elements" not in layout:
        raise ValueError(f"layout.json 格式无效：缺少 elements 字段（{layout_path}）")
    return layout


def build_pptx_bytes(layout: dict, images: Mapping[str, bytes], scale: float = 1.0) -> io.BytesIO:
    """核心：按 layout 与图片字节生成 .pptx，返回内存中的 BytesIO。

    images 为 { 元素文件名: PNG 字节 } 的映射，缺图时跳过并告警。
    """
    # 画布尺寸：优先用 layout 记录的像素宽高，缺省回退到标准 16:9；
    # python-pptx 要求边长在 1~56 英寸之间，过小的画布钳制到下限。
    width_in = max(1.0, layout.get("width", DEFAULT_SLIDE_W * PX_PER_INCH) / PX_PER_INCH * scale)
    height_in = max(1.0, layout.get("height", DEFAULT_SLIDE_H * PX_PER_INCH) / PX_PER_INCH * scale)

    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = 空白版式

    for el in layout["elements"]:
        file = el.get("file")
        data = images.get(file) if file else None
        if not data:
            print(f"[警告] 跳过缺失素材：{file}", file=sys.stderr)
            continue
        left = Inches(el["x"] / PX_PER_INCH * scale)
        top = Inches(el["y"] / PX_PER_INCH * scale)
        width = Inches(el["width"] / PX_PER_INCH * scale)
        height = Inches(el["height"] / PX_PER_INCH * scale)
        slide.shapes.add_picture(io.BytesIO(data), left, top, width, height)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_ppt(layout_path: Path, output_path: Path, scale: float = 1.0) -> Path:
    """CLI 入口：从磁盘读取 layout.json 与素材图片，写出 .pptx。"""
    layout = load_layout(layout_path)
    images: dict[str, bytes] = {}
    for el in layout["elements"]:
        file = el.get("file")
        if not file:
            continue
        img_path = layout_path.parent / file
        if img_path.is_file():
            images[file] = img_path.read_bytes()

    output_path.write_bytes(build_pptx_bytes(layout, images, scale).getvalue())
    print(f"已生成 {output_path}（{len(layout['elements'])} 个元素）")
    return output_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="读取 SlideCutter 导出的 layout.json，自动生成 PPT"
    )
    parser.add_argument("layout", type=Path, help="layout.json 路径（或其所在目录）")
    parser.add_argument(
        "-o", "--output", type=Path, default=None,
        help="输出 .pptx 路径（默认与 layout.json 同目录）",
    )
    parser.add_argument(
        "-s", "--scale", type=float, default=1.0,
        help="整体缩放比例（默认 1.0，按 96 DPI 原尺寸摆放）",
    )
    args = parser.parse_args(argv)

    layout_path = args.layout if args.layout.is_file() else args.layout / "layout.json"
    output = args.output or layout_path.with_name(
        f"{layout_path.stem}-slide.pptx"
    )
    try:
        build_ppt(layout_path, output, args.scale)
    except (FileNotFoundError, ValueError) as exc:
        print(f"[错误] {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
